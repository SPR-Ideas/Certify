"""This scripts is used to to generate Certificates."""

import subprocess
import os
import sys
import time
from shutil import rmtree
from pptx import Presentation
import pandas
# import eel


PPT_DIR = "ppt/"
CERTIFCATE_DIR = ""
MAPPED_LOG = {} # It stores the filename and username
# MAPPED_VARIABLES = {}
LAST_VAL = 0

def error(msg):
    """Print the message and exits the program."""
    print(msg)
    sys.exit(1)


def ppt_to_pdf(source_dir ,destination_dir , rm_src=True):
    """
    Args:
        source_dir (str) : Directory in which the .pptx files are located.
        destination_dir (str) : Directory in which the .pdf files should be stored.

        rm_scr (bool) : It removes the source_dir if it is set to True after conveting
                        all the .pptx files to pdf files. By default it is set to False.

    Description:
        It convertes the all the .pptx files from the source_dir to .pdf format and removes
        the source_dir if the rm_scr is set to True.

    """
    print("************ Started Creating Certificates *********")

    process = subprocess.Popen(
            f'lowriter --convert-to pdf {source_dir}*.pptx --outdir {destination_dir}',
            shell=True,stdout=subprocess.PIPE
            )

    for line in process.stdout:
        print(line.decode())
        # eel.display_log(line.decode())

    # process.wait() # It waits till the above process get completed.
    # show_progression()

    if rm_src:
        rmtree(PPT_DIR)

    print("======================X Done X===================")



def search_and_replace(find_and_replace, input_file, output_file):
    """"Search and replace text in PowerPoint while preserving formatting."""
    global MAPPED_LOG
    prs = Presentation(input_file)

    try :
        slide = prs.slides[0]
    except IndexError:
        error("Corrupted Template file {}".format(input_file))

    for shape in slide.shapes:
        if shape.has_text_frame:
            for search_str , repl_str in find_and_replace.items():

                if search_str == "{{email}}":
                    MAPPED_LOG[repl_str] = CERTIFCATE_DIR+output_file.replace("pptx","pdf")

                if(shape.text.find(search_str))!=-1:

                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text

    prs.save(PPT_DIR+output_file)



def get_details(file_name,column):
    """
    Args:
        file_name (str) : Name of the file in which a details are present.
        column str() : Name of the column in csv file to be read.

    Description:
        It Gets the details of from the file.
    """

    # Reads the given colom and stores as pydantic data_frame
    if os.path.exists(file_name):
        data = pandas.read_csv(
            file_name,
            usecols=column,
            )
    else:
        error(f"{file_name} does not exist.")

    # Returs non-duplicates dataset in reference with Mail-ID.
    return data


def create_directories(certificates,ppt_files):
    """
        It creates the required directory.
    """
    global CERTIFCATE_DIR, PPT_DIR
    if certificates[-1] != '/':
        certificates+="/"
    CERTIFCATE_DIR = certificates
    PPT_DIR = ppt_files

    os.makedirs(certificates,exist_ok=True)
    os.makedirs(ppt_files)


def generate_certificates(template,csv_file,mapped_values,certificates):
    """
    Args:
        template (str) : The file location of the Template object.
        csv_file (str) : The file location of the CSV File.

    Description:
        It Generates certificates for the members in csv file in
        accordance with template file.
    """

    create_directories(certificates,PPT_DIR)

    print("************ Intializing the Environment ***********")
    print(f"\n[*] Got Template file {template}")

    details = get_details(csv_file,list(mapped_values.values()))
    print(f"[*] Read csv file {csv_file}\n")

    # Itreate to Name list and creates a file with .pptx extension.
    for csv_iterator in details.iterrows():

        find_replace = {

                key:csv_iterator[1][values] for key,values in mapped_values.items()
            }

        search_and_replace(
            find_replace,
            template,                   # File location of the template.
            "{}.pptx".format(           # FileName of the output .pptx file.
                "certify-"
                + str(time.time()).replace(".","")   # certify-12211212.pptx
            )
        )
    print("saving to : ",CERTIFCATE_DIR)
    ppt_to_pdf(PPT_DIR,CERTIFCATE_DIR)


def start(templates,csv_file,mapped_values,save_dir):
    """
        It starts the Script.
    """
    # global MAPPED_VARIABLES
    # MAPPED_VARIABLES = mapped_values
    # st = time.time()
    generate_certificates(
        templates,
        csv_file,
        mapped_values,
        save_dir
        )
    # print(time.time()-st)
    print(MAPPED_LOG)


def progression_bar(percentage):
    """
        Just prints the progression bar with specified percentage.
    """
    space_count = 20
    progression=""
    progression = "#"*int((percentage/100 *20))
    print("[{}]".format(progression.ljust(space_count)),end='\r')


def ls_dir(dir):
    """
        Args:
            dir (str) : folder path.
        Description:
            Returns the all files and sub directory present in
            the mentioned folder.
    """
    if not os.path.exists(dir):
        error("path does not exist")

    return os.listdir(dir)


def get_precentage(src_dir_files,des_dir_files):
    """
        Args:
            src_dir_files (list) : list of files in src directory
            des_dir_files (list) : list of files in des directory
        Description:
            It checks the total files present in both the directory
            and returns a percentage of completion
            eg:
                src-files : 60 and des-files :45
                it returns 75%
    """
    global LAST_VAL
    total = len(src_dir_files)
    completed = len(des_dir_files)
    percentage = int(completed *100/total)

    if not LAST_VAL :
        LAST_VAL = percentage

    if percentage > LAST_VAL:
        LAST_VAL = percentage
    else:
        percentage =LAST_VAL

    return percentage


def show_progression():
    """
        It is used for showing progression by calling the
        progression_bar() interative manner with dynamic percentage
        as argument.
    """
    src_list = ls_dir(PPT_DIR)
    des_list =  ls_dir(CERTIFCATE_DIR)
    while src_list < des_list:
        time.sleep(0.2)
        percentage = get_precentage(src_list,des_list)
        progression_bar(percentage)
        des_list = ls_dir(CERTIFCATE_DIR)


# if __name__ == "__main__":
#     start(
#         "templates/templates.pptx",
#         "workshop.csv",
#         {"{{Full Name}}":"Name",},
#         "certificates/"
#     )
